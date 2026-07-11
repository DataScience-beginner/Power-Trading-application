"""Generate the editable Innowatt competitor and AI strategy presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "docs" / "presentations" / "innowatt_competitor_ai_roadmap.pptx"

NAVY = RGBColor(10, 28, 52)
BLUE = RGBColor(34, 102, 238)
CYAN = RGBColor(27, 181, 210)
GREEN = RGBColor(20, 160, 120)
AMBER = RGBColor(241, 166, 46)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(241, 246, 252)
MUTED = RGBColor(93, 111, 132)
DARK = RGBColor(22, 34, 51)


def add_background(slide, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_text(slide, text, x, y, w, h, size=20, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.65, 0.35, 12, 0.55, 27, NAVY, True)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.02), Inches(1.25), Inches(0.06)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BLUE
    slide.shapes[-1].line.fill.background()
    if subtitle:
        add_text(slide, subtitle, 2.05, 0.83, 10.6, 0.35, 11, MUTED)


def add_footer(slide, number):
    add_text(slide, "INNOWATT ENERGY AI  |  CONFIDENTIAL STRATEGY", 0.65, 7.15, 7, 0.2, 8, MUTED)
    add_text(slide, str(number), 12.05, 7.1, 0.55, 0.25, 9, MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, title, body, x, y, w, h, accent=BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(215, 226, 239)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(slide, title, x + 0.2, y + 0.15, w - 0.35, 0.35, 15, NAVY, True)
    add_text(slide, body, x + 0.2, y + 0.58, w - 0.38, h - 0.68, 10.5, DARK)


def add_bullets(slide, items, x, y, w, h, size=16, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
    return box


def standard_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, title, subtitle)
    add_footer(slide, len(prs.slides))
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)
    add_text(slide, "INNOWATT", 0.75, 0.55, 4, 0.5, 17, CYAN, True)
    add_text(slide, "Competitor Intelligence\n& AI Execution Roadmap", 0.75, 1.55, 8.8, 1.65, 34, WHITE, True)
    add_text(slide, "From energy dashboards to an explainable, agentic power-procurement operating system", 0.78, 3.55, 7.9, 0.8, 17, RGBColor(207, 221, 239))
    add_text(slide, "Solar  •  IEX  •  TNEB  •  15-minute intelligence", 0.78, 5.65, 6.5, 0.4, 14, CYAN, True)
    add_text(slide, "Strategy baseline — 11 July 2026", 0.78, 6.55, 4.5, 0.3, 10, RGBColor(167, 185, 208))

    slide = standard_slide(prs, "Executive decision", "What the competitor benchmark means for Innowatt")
    add_card(slide, "The market pattern", "Leaders combine governed energy data, domain forecasts, constraint-aware optimization, controlled automation, and audit-ready explanations.", 0.75, 1.45, 3.8, 2.0, BLUE)
    add_card(slide, "Our positioning", "AI-assisted procurement for Indian C&I clients: forecast demand, solar and IEX; optimize Solar/IEX/TNEB; explain every recommendation.", 4.78, 1.45, 3.8, 2.0, GREEN)
    add_card(slide, "First product move", "Launch read-only Data Quality and Market Explanation capabilities before forecasting or autonomous trading.", 8.8, 1.45, 3.8, 2.0, AMBER)
    add_text(slide, "Generic chatbot", 1.2, 4.4, 2.3, 0.4, 16, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "→", 3.55, 4.35, 0.5, 0.5, 24, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, "Domain intelligence", 4.2, 4.4, 2.6, 0.4, 16, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, "→", 6.95, 4.35, 0.5, 0.5, 24, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, "Procurement decision", 7.6, 4.4, 2.7, 0.4, 16, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, "→", 10.4, 4.35, 0.5, 0.5, 24, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, "Controlled action", 10.85, 4.4, 1.8, 0.4, 16, GREEN, True, PP_ALIGN.CENTER)

    slide = standard_slide(prs, "Competitive landscape", "Ten official product platforms reviewed")
    competitors = [
        ("GridBeyond", "Robotic trading & flexibility"), ("Fluence", "Probabilistic bidding"),
        ("Stem", "Storage dispatch"), ("Kraken", "Utility operating system"),
        ("Gridmatic", "Automated market bids"), ("Pexapark", "PPA risk & pricing"),
        ("Bidgely", "Meter AI & agents"), ("Siemens", "Digital grid twin"),
        ("Schneider", "Enterprise energy data"), ("EnergyCAP", "Bills & anomalies"),
    ]
    for i, (name, feature) in enumerate(competitors):
        col, row = i % 5, i // 5
        add_card(slide, name, feature, 0.55 + col * 2.55, 1.55 + row * 2.25, 2.25, 1.55, [BLUE, CYAN, GREEN, AMBER, BLUE][col])
    add_text(slide, "Common denominator: data → forecast → optimize → explain → automate with controls", 1.2, 6.2, 10.9, 0.45, 17, NAVY, True, PP_ALIGN.CENTER)

    slide = standard_slide(prs, "Trading and optimization leaders", "What their capabilities solve")
    add_card(slide, "GridBeyond", "70+ data sources, ML forecasts, solvers and robotic trading monetize flexible assets.", 0.75, 1.45, 3.75, 1.65, BLUE)
    add_card(slide, "Fluence Mosaic", "Probabilistic price forecasts and stochastic optimization prepare constraint-compliant bids.", 4.78, 1.45, 3.75, 1.65, CYAN)
    add_card(slide, "Gridmatic", "Weather, grid and market models drive automated price forecasts and market bids.", 8.8, 1.45, 3.75, 1.65, GREEN)
    add_card(slide, "Stem PowerTrack", "Forecasting, value stacking and real-time dispatch optimize storage and hybrid assets.", 0.75, 3.55, 3.75, 1.65, AMBER)
    add_card(slide, "Pexapark", "Forward curves and risk-adjusted valuation expose PPA and renewable capture-price risk.", 4.78, 3.55, 3.75, 1.65, BLUE)
    add_card(slide, "Lesson for Innowatt", "A recommendation needs uncertainty, constraints, risk tolerance, scenarios and reproducibility.", 8.8, 3.55, 3.75, 1.65, GREEN)

    slide = standard_slide(prs, "Enterprise data and operations leaders", "The operating foundation behind useful AI")
    add_card(slide, "Kraken", "Unified customer, billing, meter and DER workflows through a modular operating platform.", 0.75, 1.5, 3.75, 1.75, BLUE)
    add_card(slide, "Bidgely", "Energy-specific ML, disaggregation, APIs/MCP and governed agent workflows.", 4.78, 1.5, 3.75, 1.75, CYAN)
    add_card(slide, "Siemens", "Shared digital twin, data validation, modular analytics and open APIs.", 8.8, 1.5, 3.75, 1.75, GREEN)
    add_card(slide, "Schneider", "Centralized bills, intervals, procurement, risk and sustainability intelligence.", 2.75, 3.75, 3.75, 1.75, AMBER)
    add_card(slide, "EnergyCAP", "Normalized portfolio data, automated bill validation, anomalies and finance reporting.", 6.82, 3.75, 3.75, 1.75, BLUE)

    slide = standard_slide(prs, "Innowatt differentiation", "Designed for Indian C&I procurement")
    items = [
        "15-minute block-level planning and traceability",
        "DAM, RTM, GDAM and energy-schedule ingestion",
        "Monthly solar banking balance and expiry awareness",
        "Solar vs IEX vs TNEB cost and risk optimization",
        "Client and portfolio isolation for SaaS operation",
        "Explainable recommendations with approval controls",
        "Autonomous BI/data operations with complete audit trails",
    ]
    add_bullets(slide, items, 0.9, 1.45, 6.2, 4.9, 16)
    add_card(slide, "Product promise", "Forecast demand, solar availability and IEX prices—then recommend the lowest-risk, lowest-cost procurement mix with evidence.", 7.45, 1.65, 4.85, 2.1, GREEN)
    add_card(slide, "Defensible advantage", "Localized market rules + client history + trusted data + optimization outcomes create a compounding domain-data moat.", 7.45, 4.1, 4.85, 1.8, BLUE)

    slide = standard_slide(prs, "Target AI architecture", "Generative AI explains and orchestrates; domain services decide")
    layers = [
        ("Sources", "Files • Email • Meters • Weather • IEX • Tariffs", MUTED),
        ("Trust", "Ingestion • Validation • Lineage • Tenant scope", BLUE),
        ("Intelligence", "Demand • Solar • Price forecasts", CYAN),
        ("Decision", "Constraint-aware procurement optimizer", GREEN),
        ("Experience", "Dashboard • Chat • Alerts • Approval", AMBER),
        ("Action", "Controlled execution • Reconciliation • Audit", NAVY),
    ]
    for i, (name, body, color) in enumerate(layers):
        y = 1.3 + i * 0.86
        block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(y), Inches(10.9), Inches(0.62))
        block.fill.solid(); block.fill.fore_color.rgb = color; block.line.fill.background()
        add_text(slide, name, 1.45, y + 0.12, 1.7, 0.3, 13, WHITE, True)
        add_text(slide, body, 3.2, y + 0.12, 8.4, 0.3, 13, WHITE)

    slide = standard_slide(prs, "AI integration roadmap", "Value and autonomy increase only after trust gates")
    phases = [
        ("AI-0", "Trust foundation", "Contracts, lineage, audit, tenant authorization"),
        ("AI-1", "Read-only intelligence", "Data quality, explanations, evidence-backed chat"),
        ("AI-2", "Forecasting", "Demand, solar and IEX forecasts with quantiles"),
        ("AI-3", "Optimization", "Solar/IEX/TNEB scenarios and recommendations"),
        ("AI-4", "Controlled autonomy", "Email-to-dashboard, approvals and bounded actions"),
    ]
    for i, (phase, name, body) in enumerate(phases):
        x = 0.55 + i * 2.55
        add_card(slide, f"{phase}  {name}", body, x, 1.65, 2.25, 3.35, [MUTED, BLUE, CYAN, GREEN, AMBER][i])
        if i < 4:
            add_text(slide, "→", x + 2.27, 3.0, 0.25, 0.35, 18, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, "Gate each phase on evidence quality, tenant isolation, accuracy, reproducibility and audit completeness", 1.0, 5.65, 11.4, 0.6, 16, NAVY, True, PP_ALIGN.CENTER)

    slide = standard_slide(prs, "Agentic way of working", "Specialists govern; agents execute bounded repeatable workflows")
    flow = ["Trigger", "Classify", "Validate", "Plan", "Use tool", "Verify", "Audit", "Escalate"]
    for i, step in enumerate(flow):
        x = 0.35 + i * 1.61
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(1.35), Inches(0.72))
        shape.fill.solid(); shape.fill.fore_color.rgb = BLUE if i < 6 else GREEN; shape.line.fill.background()
        add_text(slide, step, x + 0.05, 1.72, 1.25, 0.25, 12, WHITE, True, PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            add_text(slide, "→", x + 1.37, 1.7, 0.22, 0.25, 13, MUTED, True, PP_ALIGN.CENTER)
    add_card(slide, "Agents", "Email intake • file classification • data quality • parser/upload • dashboard analyst • forecast monitor • procurement advisor • notification", 0.8, 3.0, 5.7, 2.3, CYAN)
    add_card(slide, "Specialists", "Define policy • approve consequential actions • resolve exceptions • improve models and strategy • own client relationships", 6.85, 3.0, 5.7, 2.3, AMBER)
    add_text(slide, "Autonomy ladder: Shadow → Assist → Supervised automation → Bounded autonomy", 1.1, 5.85, 11.1, 0.45, 16, NAVY, True, PP_ALIGN.CENTER)

    slide = standard_slide(prs, "Agent control plane", "How hundreds or thousands of agents stay manageable")
    add_card(slide, "Registry", "Identity • owner • version • status • triggers • tools • model policy", 0.75, 1.45, 3.75, 1.7, BLUE)
    add_card(slide, "Permissions", "Tenant scope • read/write bounds • approvals • budgets • kill switch", 4.78, 1.45, 3.75, 1.7, GREEN)
    add_card(slide, "Observability", "Runs • tool calls • evidence • confidence • latency • cost • outcomes", 8.8, 1.45, 3.75, 1.7, CYAN)
    add_card(slide, "Evaluation", "Golden datasets • policy tests • replay • regression • drift detection", 2.75, 3.75, 3.75, 1.7, AMBER)
    add_card(slide, "Lifecycle", "Proposed → active → deprecated → retired, with rollback and ownership", 6.82, 3.75, 3.75, 1.7, BLUE)

    slide = standard_slide(prs, "Technology stack", "Build on the current product; add components only when justified")
    add_card(slide, "Frontend", "React + TypeScript\nMaterial UI\nTanStack Query\nForms + schema validation\nCharts, evidence, confidence, approval UI\nSSE for agent progress", 0.65, 1.35, 3.85, 4.8, BLUE)
    add_card(slide, "Backend & data", "FastAPI + Pydantic\nSQLAlchemy + PostgreSQL + Alembic\nObject storage for raw files\nWorker queue + Redis when needed\nOpenTelemetry and structured audit", 4.75, 1.35, 3.85, 4.8, GREEN)
    add_card(slide, "AI, ML & optimization", "Provider-neutral AI adapter\nVersioned prompts and evaluations\nMLflow/equivalent model registry\nBaseline ML before complex models\nOR-Tools or Pyomo optimizer\nTenant-scoped retrieval and tools", 8.85, 1.35, 3.85, 4.8, CYAN)

    slide = standard_slide(prs, "Phase-wise execution", "Indicative ranges—re-estimate after foundation discovery")
    rows = [
        ("0", "2–3 weeks", "Contracts, lineage, audit, tenant authorization"),
        ("1", "3–4 weeks", "Data Quality Agent, explanations, read-only chatbot"),
        ("2", "6–8 weeks", "Demand, solar and next-day IEX forecasting"),
        ("3", "6–8 weeks", "Procurement optimizer and scenario backtests"),
        ("4", "6–10 weeks", "Email-to-dashboard and approval workflows"),
        ("5", "Governance gated", "Controlled market submission and reconciliation"),
    ]
    for i, (phase, duration, output) in enumerate(rows):
        y = 1.35 + i * 0.82
        color = [MUTED, BLUE, CYAN, GREEN, AMBER, NAVY][i]
        add_text(slide, f"PHASE {phase}", 0.85, y + 0.12, 1.3, 0.28, 12, color, True)
        add_text(slide, duration, 2.35, y + 0.12, 1.55, 0.28, 12, NAVY, True)
        add_text(slide, output, 4.05, y + 0.12, 7.9, 0.32, 13, DARK)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(y + 0.63), Inches(11.4), Inches(0.015))
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(220, 229, 239); line.line.fill.background()

    slide = standard_slide(prs, "First implementation epic", "Data Quality + Market Explanation")
    add_bullets(slide, [
        "Define typed Pydantic request and response contracts",
        "Run deterministic quality rules before any LLM explanation",
        "Retrieve evidence only within client, portfolio and date scope",
        "Return finding, severity, evidence, confidence and recommended next step",
        "Create market explanations from approved analytics APIs",
        "Record model/prompt version, inputs, outputs, feedback and latency",
        "Add UI evidence panels and rigorous authorization/evaluation tests",
    ], 0.85, 1.35, 6.4, 5.2, 15)
    add_card(slide, "Release boundary", "Read-only in production. No automatic correction, schedule mutation, bid placement or external notification.", 7.6, 1.55, 4.65, 1.65, AMBER)
    add_card(slide, "Success", "≥90% seeded defect detection • zero cross-client exposure • evidence on every answer • useful specialist feedback", 7.6, 3.65, 4.65, 1.85, GREEN)

    slide = standard_slide(prs, "Governance and success measures", "AI earns autonomy through measured reliability")
    add_card(slide, "Data", "Completeness • freshness • schema validity • lineage coverage", 0.75, 1.45, 3.75, 1.6, BLUE)
    add_card(slide, "Models", "Baseline lift • calibration • drift • client/block accuracy", 4.78, 1.45, 3.75, 1.6, CYAN)
    add_card(slide, "Optimization", "Feasibility • savings • risk • replay reproducibility", 8.8, 1.45, 3.75, 1.6, GREEN)
    add_card(slide, "Agents", "Completion • exception rate • false actions • cost • latency", 2.75, 3.65, 3.75, 1.6, AMBER)
    add_card(slide, "Business", "TNEB exposure • solar expiry • procurement cost • analyst time", 6.82, 3.65, 3.75, 1.6, BLUE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)
    add_text(slide, "Decision requested", 0.8, 0.8, 5.5, 0.55, 18, CYAN, True)
    add_text(slide, "Start Phase AI-0 + AI-1", 0.8, 1.6, 8.2, 0.8, 34, WHITE, True)
    add_text(slide, "Build the trusted Data Quality and Market Explanation foundation before predictive procurement and autonomous workflows.", 0.82, 2.8, 9.6, 0.95, 19, RGBColor(207, 221, 239))
    add_card(slide, "Immediate sprint", "Contracts • quality rules • tenant-safe evidence • audit events • read-only APIs • dashboard explanation UI • evaluation fixtures", 0.82, 4.45, 7.2, 1.45, GREEN)
    add_text(slide, "Innowatt Energy AI", 9.2, 6.55, 3.25, 0.35, 14, CYAN, True, PP_ALIGN.RIGHT)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Generated {OUTPUT}")


if __name__ == "__main__":
    build_deck()
